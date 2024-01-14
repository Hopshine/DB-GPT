from typing import Any, List, Optional

from dbgpt.rag.chunk import Document
from dbgpt.rag.knowledge.base import (
    ChunkStrategy,
    DocumentType,
    Knowledge,
    KnowledgeType,
)


class PPTXKnowledge(Knowledge):
    """PPTX Knowledge"""

    def __init__(
        self,
        file_path: Optional[str] = None,
        knowledge_type: KnowledgeType = KnowledgeType.DOCUMENT,
        loader: Optional = None,
        language: Optional[str] = "zh",
        **kwargs: Any,
    ) -> None:
        """Initialize with PDF Knowledge arguments.
        Args:
            file_path:(Optional[str]) file path
            knowledge_type:(KnowledgeType) knowledge type
            loader:(Optional[Any]) loader
        """
        self._path = file_path
        self._type = knowledge_type
        self._loader = loader
        self._language = language

    def _load(self) -> List[Document]:
        """Load pdf document from loader"""
        if self._loader:
            documents = self._loader.load()
        else:
            from pptx import Presentation

            pr = Presentation(self._path)
            docs = []
            for slide in pr.slides:
                content = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        content += shape.text
                docs.append(
                    Document(content=content, metadata={"source": slide.slide_id})
                )
            return docs
        return [Document.langchain2doc(lc_document) for lc_document in documents]

    @classmethod
    def support_chunk_strategy(cls) -> List[ChunkStrategy]:
        return [
            ChunkStrategy.CHUNK_BY_SIZE,
            ChunkStrategy.CHUNK_BY_PAGE,
            ChunkStrategy.CHUNK_BY_SEPARATOR,
        ]

    @classmethod
    def default_chunk_strategy(cls) -> ChunkStrategy:
        return ChunkStrategy.CHUNK_BY_SIZE

    @classmethod
    def type(cls) -> KnowledgeType:
        return KnowledgeType.DOCUMENT

    @classmethod
    def document_type(cls) -> DocumentType:
        return DocumentType.PPTX
